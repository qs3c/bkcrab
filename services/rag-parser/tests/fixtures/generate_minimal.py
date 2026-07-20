from __future__ import annotations

import argparse
import zipfile
from pathlib import Path

from PIL import Image as PILImage


def _sample_image(path: Path) -> None:
    image = PILImage.new("RGB", (24, 16), color=(35, 99, 180))
    image.save(path, format="PNG")


def _docx(path: Path, image_path: Path) -> None:
    content_types = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Default Extension="png" ContentType="image/png"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    root_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    document_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rIdImage1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/image1.png"/>
</Relationships>"""
    document = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:pic="http://schemas.openxmlformats.org/drawingml/2006/picture">
 <w:body>
  <w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>Minimal DOCX</w:t></w:r></w:p>
  <w:p><w:r><w:t>Before the image</w:t></w:r></w:p>
  <w:p><w:r><w:drawing><wp:inline>
   <wp:extent cx="914400" cy="609600"/>
   <wp:docPr id="1" name="Diagram" descr="Blue architecture diagram"/>
   <a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">
    <pic:pic><pic:nvPicPr><pic:cNvPr id="1" name="sample.png"/><pic:cNvPicPr/></pic:nvPicPr>
     <pic:blipFill><a:blip r:embed="rIdImage1"/><a:stretch><a:fillRect/></a:stretch></pic:blipFill>
     <pic:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="914400" cy="609600"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></pic:spPr>
    </pic:pic>
   </a:graphicData></a:graphic>
  </wp:inline></w:drawing></w:r></w:p>
  <w:p><w:r><w:t>After the image</w:t></w:r></w:p>
  <w:sectPr/>
 </w:body>
</w:document>"""
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", root_rels)
        archive.writestr("word/document.xml", document)
        archive.writestr("word/_rels/document.xml.rels", document_rels)
        archive.write(image_path, "word/media/image1.png")


def _pptx(path: Path, image_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(5), Inches(0.5))
    textbox.text_frame.text = "Minimal PPTX"
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1.2), width=Inches(2)
    )
    picture._element.nvPicPr.cNvPr.set("descr", "Blue slide diagram")
    presentation.save(path)


def _xlsx(path: Path, image_path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet.append(["Name", "Value"])
    sheet.append(["alpha", 1])
    picture = Image(str(image_path))
    picture.anchor = "B3"
    sheet.add_image(picture)
    workbook.save(path)


def generate_all(output: Path) -> dict[str, Path]:
    output.mkdir(parents=True, exist_ok=True)
    image_path = output / "sample.png"
    _sample_image(image_path)
    paths = {
        "docx": output / "minimal.docx",
        "pptx": output / "minimal.pptx",
        "xlsx": output / "minimal.xlsx",
    }
    _docx(paths["docx"], image_path)
    _pptx(paths["pptx"], image_path)
    _xlsx(paths["xlsx"], image_path)
    return paths


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate legal minimal Office fixtures")
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    generated = generate_all(args.output)
    for source_format, path in generated.items():
        print(f"{source_format}: {path}")


if __name__ == "__main__":
    main()
