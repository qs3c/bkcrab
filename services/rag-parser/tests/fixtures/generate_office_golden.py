from __future__ import annotations

import argparse
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image as PILImage

PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _sample_image(path: Path) -> None:
    image = PILImage.new("RGB", (40, 24), color=(35, 99, 180))
    image.save(path, format="PNG")


def _docx(path: Path, image_path: Path) -> None:
    content_types = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Default Extension="png" ContentType="image/png"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
  <Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>
  <Override PartName="/word/numbering.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.numbering+xml"/>
</Types>"""
    root_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    document_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rIdImage1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/image1.png"/>
  <Relationship Id="rIdStyles" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
  <Relationship Id="rIdNumbering" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/numbering" Target="numbering.xml"/>
</Relationships>"""
    styles = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:style w:type="paragraph" w:default="1" w:styleId="Normal"><w:name w:val="Normal"/></w:style>
  <w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/><w:basedOn w:val="Normal"/><w:pPr><w:outlineLvl w:val="0"/></w:pPr></w:style>
  <w:style w:type="paragraph" w:styleId="Code"><w:name w:val="Code"/><w:basedOn w:val="Normal"/></w:style>
  <w:style w:type="paragraph" w:styleId="Preformatted"><w:name w:val="Preformatted"/><w:basedOn w:val="Normal"/></w:style>
</w:styles>"""
    numbering = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:numbering xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:abstractNum w:abstractNumId="0"><w:multiLevelType w:val="singleLevel"/><w:lvl w:ilvl="0"><w:numFmt w:val="bullet"/><w:lvlText w:val="&#x2022;"/></w:lvl></w:abstractNum>
  <w:num w:numId="1"><w:abstractNumId w:val="0"/></w:num>
</w:numbering>"""

    drawing = """<w:p><w:r><w:drawing><wp:inline>
 <wp:extent cx="1143000" cy="685800"/>
 <wp:docPr id="{drawing_id}" name="Architecture" descr="Architecture diagram"/>
 <a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">
  <pic:pic><pic:nvPicPr><pic:cNvPr id="{drawing_id}" name="office-golden.png"/><pic:cNvPicPr/></pic:nvPicPr>
   <pic:blipFill><a:blip r:embed="rIdImage1"/><a:stretch><a:fillRect/></a:stretch></pic:blipFill>
   <pic:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="1143000" cy="685800"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></pic:spPr>
  </pic:pic>
 </a:graphicData></a:graphic>
</wp:inline></w:drawing></w:r></w:p>"""
    document = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:pic="http://schemas.openxmlformats.org/drawingml/2006/picture">
 <w:body>
  <w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>Office Golden</w:t></w:r></w:p>
  <w:p><w:r><w:t>Intro before list.</w:t></w:r></w:p>
  <w:p><w:pPr><w:numPr><w:ilvl w:val="0"/><w:numId w:val="1"/></w:numPr></w:pPr><w:r><w:t>First item</w:t></w:r></w:p>
  <w:p><w:pPr><w:numPr><w:ilvl w:val="0"/><w:numId w:val="1"/></w:numPr></w:pPr><w:r><w:t>Second item</w:t></w:r></w:p>
  <w:tbl><w:tblPr/><w:tblGrid><w:gridCol/><w:gridCol/></w:tblGrid>
   <w:tr><w:tc><w:p><w:r><w:t>Product</w:t></w:r></w:p></w:tc><w:tc><w:p><w:r><w:t>Count</w:t></w:r></w:p></w:tc></w:tr>
   <w:tr><w:tc><w:p><w:r><w:t>Crab</w:t></w:r></w:p></w:tc><w:tc><w:p><w:r><w:t>2</w:t></w:r></w:p></w:tc></w:tr>
  </w:tbl>
  <w:p><w:pPr><w:pStyle w:val="Code"/></w:pPr><w:r><w:t>print('alpha')</w:t></w:r></w:p>
  <w:p><w:pPr><w:pStyle w:val="Preformatted"/></w:pPr><w:r><w:t>SELECT * FROM docs;</w:t></w:r></w:p>
  <w:p><w:r><w:rPr><w:rFonts w:ascii="Courier New" w:hAnsi="Courier New"/></w:rPr><w:t>ordinary monospace stays prose</w:t></w:r></w:p>
  <w:p><w:r><w:t>Before image.</w:t></w:r></w:p>
  {drawing.format(drawing_id=1)}
  <w:p><w:r><w:t>Between repeated images.</w:t></w:r></w:p>
  {drawing.format(drawing_id=2)}
  <w:p><w:r><w:t>After image.</w:t></w:r></w:p>
  <w:sectPr/>
 </w:body>
</w:document>""".encode()
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", root_rels)
        archive.writestr("word/document.xml", document)
        archive.writestr("word/styles.xml", styles)
        archive.writestr("word/numbering.xml", numbering)
        archive.writestr("word/_rels/document.xml.rels", document_rels)
        archive.write(image_path, "word/media/image1.png")


def _set_picture_alt(picture, value: str) -> None:
    picture._element.nvPicPr.cNvPr.set("descr", value)


def _pptx(path: Path, image_path: Path) -> None:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    presentation = Presentation()
    rich_slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    bottom = rich_slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(5), Inches(0.4)
    )
    bottom.text_frame.text = "Bottom text"

    right_picture = rich_slide.shapes.add_picture(
        str(image_path), Inches(4), Inches(1.5), width=Inches(1.25)
    )
    _set_picture_alt(right_picture, "Architecture diagram")
    left_picture = rich_slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(1.5), width=Inches(1.25)
    )
    _set_picture_alt(left_picture, "Architecture diagram")

    group = rich_slide.shapes.add_group_shape()
    group_right = group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.4), Inches(1.5), Inches(0.5)
    )
    group_right.text = "Group right"
    group_left = group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.4), Inches(1.5), Inches(0.5)
    )
    group_left.text = "Group left"

    code = rich_slide.shapes.add_textbox(
        Inches(0.8), Inches(4.4), Inches(4), Inches(0.6)
    )
    code.name = "Code"
    code.text_frame.text = "def slide_code():\n    return 7"

    top = rich_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(5), Inches(0.4)
    )
    top.text_frame.text = "Top text"
    rich_slide.notes_slide.notes_text_frame.text = "Remember the architecture caveat."

    second_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_text = second_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(5), Inches(0.4)
    )
    second_text.text_frame.text = "Presentation-order first slide"
    presentation.save(path)

    # Deliberately make presentation order differ from slide part numbering.
    # Consumers must follow p:sldIdLst relationships, not fuzzy slide filenames.
    with zipfile.ZipFile(path) as source:
        values = {info.filename: source.read(info) for info in source.infolist()}
    root = ET.fromstring(values["ppt/presentation.xml"])
    slide_ids = root.find(f"{{{PRESENTATION_NS}}}sldIdLst")
    assert slide_ids is not None and len(slide_ids) == 2
    slide_ids[:] = [slide_ids[1], slide_ids[0]]
    values["ppt/presentation.xml"] = ET.tostring(
        root, encoding="utf-8", xml_declaration=True
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as output:
        for name, value in values.items():
            output.writestr(name, value)


def _xlsx(path: Path, image_path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image

    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.append(["Name", "Value"])
    summary.append(["alpha", 1])
    summary.append(["beta", 2])
    summary_image = Image(str(image_path))
    summary_image.anchor = "B3"
    summary.add_image(summary_image)

    details = workbook.create_sheet("Details")
    details.append(["Key", "Description"])
    details.append(["k1", "detail one"])
    details_image = Image(str(image_path))
    details_image.anchor = "A2"
    details.add_image(details_image)
    workbook.save(path)

    with zipfile.ZipFile(path) as source:
        values = {info.filename: source.read(info) for info in source.infolist()}
    for name in sorted(
        item for item in values if item.startswith("xl/drawings/drawing") and item.endswith(".xml")
    ):
        root = ET.fromstring(values[name])
        for node in root.iter():
            if node.tag.rsplit("}", 1)[-1] == "cNvPr":
                node.set("descr", "Metrics chart")
        values[name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as output:
        for name, value in values.items():
            output.writestr(name, value)


def generate_all(output: Path) -> dict[str, Path]:
    output.mkdir(parents=True, exist_ok=True)
    image_path = output / "office-golden.png"
    _sample_image(image_path)
    paths = {
        "docx": output / "office-golden.docx",
        "pptx": output / "office-golden.pptx",
        "xlsx": output / "office-golden.xlsx",
    }
    _docx(paths["docx"], image_path)
    _pptx(paths["pptx"], image_path)
    _xlsx(paths["xlsx"], image_path)
    return paths


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate legal Office positioning fixtures")
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    for source_format, path in generate_all(args.output).items():
        print(f"{source_format}: {path}")


if __name__ == "__main__":
    main()
